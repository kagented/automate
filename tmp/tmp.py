import os
from os import walk
import pprint
from pptx import Presentation
import win32com.client

#filelist = ['1.ppt', '2.ppt', '3.ppt']


prs = Presentation('1.ppt')

title_slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가


prs.save('1_.ppt')

